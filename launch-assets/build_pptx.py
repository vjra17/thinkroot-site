"""
ThinkRoot Young Leaders — Parent Presentation
Generates presentation-2026-04-27.pptx from slide content.
Run with: python3 build_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Brand colours ────────────────────────────────────────────────────────────
FOREST = RGBColor(0x1a, 0x3a, 0x1a)
GOLD   = RGBColor(0xc8, 0x96, 0x3c)
CREAM  = RGBColor(0xfa, 0xf8, 0xf2)
BARK   = RGBColor(0x3d, 0x2b, 0x1a)
WHITE  = RGBColor(0xff, 0xff, 0xff)
MIST   = RGBColor(0xe8, 0xf4, 0xe8)
MID    = RGBColor(0x44, 0x55, 0x44)   # body text on light bg

# ── Slide dimensions (widescreen 16:9) ───────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]  # completely blank

# ── Helper functions ─────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, l, t, w, h, text, font_name="DM Sans", font_size=18,
                bold=False, color=FOREST, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_textbox_paras(slide, l, t, w, h, paras, font_name="DM Sans",
                      font_size=14, color=MID, line_spacing=None):
    """Add a textbox with multiple paragraphs. paras = list of (text, bold, size, color, align)"""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for (text, bold, size, clr, align) in paras:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        run.font.name = font_name
        run.font.size = Pt(size if size else font_size)
        run.font.bold = bold
        run.font.color.rgb = clr if clr else color
    return txBox

def add_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text

def header_bar(slide, slide_num, total, title):
    """Dark forest green top bar with slide number + title."""
    add_rect(slide, 0, 0, W, Inches(1.0), FOREST)
    num_txt = f"SLIDE {slide_num} OF {total}"
    add_textbox(slide, Inches(0.45), Inches(0.12), Inches(3), Inches(0.38),
                num_txt, font_size=9, color=GOLD, bold=False)
    add_textbox(slide, Inches(0.45), Inches(0.48), Inches(12), Inches(0.48),
                title, font_name="Playfair Display", font_size=20, bold=True, color=WHITE)

def cream_bg(slide):
    add_rect(slide, 0, 0, W, H, CREAM)

def gold_headline(slide, text, top=Inches(1.15), height=Inches(0.85)):
    """Gold left-border headline block."""
    add_rect(slide, Inches(0.45), top, Inches(0.06), height, GOLD)
    add_textbox(slide, Inches(0.65), top, Inches(11.8), height,
                text, font_name="Playfair Display", font_size=26,
                bold=True, color=FOREST)

def bullet_block(slide, items, top, left=Inches(0.65), width=Inches(11.8),
                 font_size=14, color=MID, bullet="–  "):
    """Stack of bullet items."""
    paras = []
    for item in items:
        paras.append((bullet + item, False, font_size, color, PP_ALIGN.LEFT))
    add_textbox_paras(slide, left, top, width, Inches(4),
                      paras, font_size=font_size, color=color)

def two_col_box(slide, left_title, left_items, right_title, right_items,
                top=Inches(2.15), col_h=Inches(3.9)):
    col_w = Inches(5.7)
    # Left card (light)
    add_rect(slide, Inches(0.45), top, col_w, col_h, MIST)
    add_textbox(slide, Inches(0.65), top + Inches(0.18), col_w - Inches(0.35), Inches(0.45),
                left_title, font_name="Playfair Display", font_size=15,
                bold=True, color=FOREST)
    paras_l = [("–  " + i, False, 13, MID, PP_ALIGN.LEFT) for i in left_items]
    add_textbox_paras(slide, Inches(0.65), top + Inches(0.65),
                      col_w - Inches(0.35), col_h - Inches(0.8), paras_l)
    # Right card (dark)
    add_rect(slide, Inches(6.45), top, col_w, col_h, FOREST)
    add_textbox(slide, Inches(6.65), top + Inches(0.18), col_w - Inches(0.35), Inches(0.45),
                right_title, font_name="Playfair Display", font_size=15,
                bold=True, color=GOLD)
    paras_r = [("–  " + i, False, 13, WHITE, PP_ALIGN.LEFT) for i in right_items]
    add_textbox_paras(slide, Inches(6.65), top + Inches(0.65),
                      col_w - Inches(0.35), col_h - Inches(0.8), paras_r,
                      color=WHITE)

def bottom_line(slide, text, top=Inches(6.4)):
    add_rect(slide, Inches(0.45), top, W - Inches(0.9), Inches(0.72), FOREST)
    add_textbox(slide, Inches(0.65), top + Inches(0.1), W - Inches(1.2), Inches(0.52),
                text, font_name="Playfair Display", font_size=13,
                bold=True, color=GOLD, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, FOREST)
# Gold accent line
add_rect(s, Inches(0.45), Inches(2.6), Inches(0.08), Inches(2.2), GOLD)
add_textbox(s, Inches(0.65), Inches(1.8), Inches(11), Inches(0.5),
            "ThinkRoot Young Leaders", font_name="DM Sans", font_size=14,
            bold=False, color=GOLD)
add_textbox(s, Inches(0.65), Inches(2.3), Inches(11), Inches(1.5),
            "Building Children\nWho Can Handle Life",
            font_name="Playfair Display", font_size=44, bold=True, color=WHITE)
add_textbox(s, Inches(0.65), Inches(4.0), Inches(11), Inches(0.45),
            "An introduction to ThinkRoot Young Leaders — Resilience Lab",
            font_name="DM Sans", font_size=16, color=MIST)
add_rect(s, Inches(0.45), Inches(6.4), W - Inches(0.9), Inches(0.001), GOLD)
add_textbox(s, Inches(0.65), Inches(6.5), Inches(11), Inches(0.45),
            "Tina Puri  |  Founder & Program Director  |  thinkyoungleaders@gmail.com",
            font_name="DM Sans", font_size=13, color=GOLD)
add_notes(s, """Welcome everyone and thank you for being here. I know your time is valuable, so I'll keep this focused, practical, and — I hope — something that stays with you long after today.""")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Wolf & Cage
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
cream_bg(s)
header_bar(s, 2, 15, "The Opening Hook: The Wolf and The Cage")
gold_headline(s, "Two animals. Two very different minds.")
two_col_box(s,
    "🐺  The Wolf",
    ["Lives in the wild. No guarantees.",
     "Stays sharp. Reads the environment.",
     "Adapts. Recovers. Persists.",
     "Reaches its full potential — because it has to."],
    "🐕  The Dog in the Cage",
    ["Comfortable. Safe. Fed on schedule.",
     "Never has to hunt. Never has to adapt.",
     "Over time — loses its instincts.",
     "Forgets how to read the environment."])
add_notes(s, """I want to start with something Warren Buffett once said about people in the workforce. He described two kinds of people — the wolf, who is always sharp because survival depends on it. And the dog in a golden cage — comfortable, provided for, but slowly losing the instincts it was born with.

He was talking about financial independence. But today I want to borrow that image for something different.

Because I believe our children's minds work exactly the same way.

A mind that is trained to handle pressure, setbacks, difficulty and uncertainty — that mind becomes sharp. It grows. It reaches its potential.

A mind that is never challenged to work with what happens inside — that mind becomes dependent. Fragile. It shuts down when things get hard.

The question I want to ask today is: which kind of mind are we building in our children?""")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Invisible Cage
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
cream_bg(s)
header_bar(s, 3, 15, "The Cage No One Talks About")
gold_headline(s, "The cage is invisible. But your child may already be living in one.")
add_textbox(s, Inches(0.65), Inches(2.1), Inches(11.5), Inches(0.38),
            "The cage is not made of bars. It is made of:", font_size=14,
            color=MID, bold=False)
bullet_block(s, [
    '"I can\'t do this."',
    "Shutting down after one wrong answer",
    "Anxiety that rises before an exam — even when they are prepared",
    "Giving up when something feels too hard",
    "Reacting big to small setbacks",
    "Confidence that only works when everything is going well"
], top=Inches(2.5))
bottom_line(s, "These are not personality traits. They are patterns. And patterns can be changed.")
add_notes(s, """The cage I'm talking about today is not made of poverty, or lack of opportunity, or lack of intelligence.

Your children have all of those things. They have loving parents who have sacrificed enormously. They have tutors, activities, the best schools, every resource.

And yet — many of them are struggling in ways that all of that cannot fix.

Because the cage is internal. It is the pattern of thought, the emotional habit, the automatic response that kicks in when life gets hard. And most children — most adults — have never been taught to work with it.""")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — You Have Given Everything
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
cream_bg(s)
header_bar(s, 4, 15, "You Have Given Your Child Everything")
gold_headline(s, "You have invested in everything that goes in.")
# Checklist items
check_items = [
    "The best academic preparation",
    "Tutoring and test prep",
    "Sports, music, extracurriculars",
    "Enrichment programs and experiences",
    "Every tool and resource they need",
]
paras = [("✓  " + item, True, 15, FOREST, PP_ALIGN.LEFT) for item in check_items]
add_textbox_paras(s, Inches(2.0), Inches(2.1), Inches(9), Inches(3.2), paras)
# Gold separator line
add_rect(s, Inches(0.65), Inches(5.5), Inches(11.5), Inches(0.002), GOLD)
add_textbox(s, Inches(0.65), Inches(5.6), Inches(11.5), Inches(0.65),
            "But no one has taught them what to do when their own mind works against them.",
            font_name="Playfair Display", font_size=18, bold=True, color=GOLD)
add_notes(s, """I want to be clear — I am not here to question your parenting. Everything on this list matters. You are doing everything right by those standards.

But there is a layer that is missing. Not from your efforts — from the entire system. From schools. From the way we were raised. From the culture we grew up in.

No one taught most of us what to do when anxiety rises before a big moment. No one taught us how to recover from failure with our self-belief intact. No one taught us how to interrupt the voice that says "I'm not good enough."

We figured it out — or we didn't. And the patterns we didn't resolve, we carry today. Our children deserve better than that.""")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — What Parents See
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
cream_bg(s)
header_bar(s, 5, 15, "What Parents in This Room Are Already Seeing")
gold_headline(s, "You have probably already seen one of these.")
scenarios = [
    '"She studies hard but panics during the exam and goes blank."',
    '"He gets one bad grade and says he\'s stupid and wants to quit."',
    '"She starts crying the moment something doesn\'t go her way."',
    '"He never tries anything new because he\'s afraid of failing."',
    '"She seems anxious all the time — I don\'t know why."',
    '"He compares himself to classmates constantly and it\'s affecting his confidence."',
]
col_h = Inches(0.72)
col_w = Inches(5.7)
for i, sc in enumerate(scenarios):
    row = i // 2
    col = i % 2
    l = Inches(0.45) + col * Inches(6.95)
    t = Inches(2.1) + row * (col_h + Inches(0.12))
    add_rect(s, l, t, col_w, col_h, MIST)
    add_textbox(s, l + Inches(0.18), t + Inches(0.08), col_w - Inches(0.3), col_h - Inches(0.12),
                "📌  " + sc, font_size=11, color=FOREST)
bottom_line(s, "This is not a discipline problem. It is a skills gap.")
add_notes(s, """I want to ask for a show of hands — but feel free to just note this privately.

How many of you have seen your child shut down after a setback? Panic before an exam they were prepared for? Give up when something got hard?

[Pause.]

Almost every hand in the room. I have seen this every time I present to parents.

This is not about a child being weak or sensitive or dramatic. This is about a child who was never given a tool for that moment. They were taught the subject. But no one taught them what to do with the fear, the frustration, the self-doubt that shows up right alongside it.""")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — The Gap
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
cream_bg(s)
header_bar(s, 6, 15, "The Gap Schools Don't Fill")
gold_headline(s, "Schools teach children what to think.\nRarely how to manage what happens within.")
two_col_box(s,
    "What schools teach:",
    ["Content, formulas, facts",
     "How to perform and achieve",
     "How to compete",
     "How to follow a system"],
    "What no school teaches:",
    ["How to recover after a mistake",
     "How to stay steady when pressure rises",
     "How to interrupt a negative thought spiral",
     "How to believe in yourself when results say otherwise",
     "How to regulate the emotional state that determines everything"],
    top=Inches(2.3))
bottom_line(s, 'There is no subject called: "What to do when your mind works against you."')
add_notes(s, """Think about your own education. Think about what you were taught.

Math, science, history, language. Study hard, work harder, perform better. The formula for achievement is very clear.

But was there a single class — a single lesson — that taught you how to handle the moment right before a big interview when your anxiety is rising? Or how to recover your confidence after a public failure? Or how to stop a thought spiral at 2am?

For most of us — no. We figured it out over decades. Many of us are still figuring it out.

Our children are entering one of the most competitive, unpredictable, high-pressure environments in history. And we are sending them in with half the toolkit.""")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Research
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
cream_bg(s)
header_bar(s, 7, 15, "What the Research Says")
gold_headline(s, "This is not a theory. The science is clear.")

research = [
    ("📊 CDC Youth Risk Behavior Survey",
     "42% of high school students report persistent sadness or hopelessness.",
     "Nearly half. This is not a fringe statistic."),
    ("📊 Harvard Business Publishing, 2024",
     "4 of the 5 most in-demand workplace skills are human skills.",
     "Self-awareness, resilience, flexibility, agility. Yet schools rarely teach them."),
    ("📊 Albert Bandura, Stanford University",
     "Self-efficacy is a stronger predictor of success than actual ability.",
     "What your child believes about themselves matters more than what they know."),
    ("📊 Yale Center for Emotional Intelligence",
     "EQ predicts adult career success more reliably than IQ alone.",
     "The skills dismissed as 'soft' are the ones that determine outcomes."),
]
rw = Inches(5.7)
rh = Inches(1.88)
for i, (src, stat, note) in enumerate(research):
    row = i // 2
    col = i % 2
    l = Inches(0.45) + col * Inches(6.95)
    t = Inches(2.1) + row * (rh + Inches(0.1))
    add_rect(s, l, t, rw, rh, FOREST)
    add_textbox(s, l + Inches(0.18), t + Inches(0.12), rw - Inches(0.3), Inches(0.38),
                src, font_size=10, bold=True, color=GOLD)
    add_textbox(s, l + Inches(0.18), t + Inches(0.5), rw - Inches(0.3), Inches(0.62),
                stat, font_size=13, bold=True, color=WHITE)
    add_textbox(s, l + Inches(0.18), t + Inches(1.18), rw - Inches(0.3), Inches(0.55),
                note, font_size=11, color=MIST)
add_notes(s, """I know this audience. You are engineers, product managers, scientists. You want data before you decide.

So let me give you the data.

[Walk through each stat briefly.]

The Stanford finding is the one I want to highlight. Self-efficacy — what does your child believe about their own ability to handle challenges?

A child with high self-efficacy tries harder, recovers faster, and achieves more — independent of their actual ability. A child with low self-efficacy gives up before they even begin.

Self-efficacy is not fixed. It is built. And it is built through exactly the kind of practice we are talking about today.""")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Why Now
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
cream_bg(s)
header_bar(s, 8, 15, "Why Childhood Is The Window")
gold_headline(s, "The best time to build these skills is now —\nbefore the patterns harden.")
points = [
    ("🧠  The brain is most malleable in childhood.",
     "Neuroplasticity research shows that habits and responses formed early are the ones that persist. The longer they go unaddressed, the harder they are to change."),
    ("⚡  Stress physically reduces cognitive function.",
     "Research from HeartMath Institute: negative emotional states measurably impair memory, focus, and decision-making. Your child cannot think clearly when in a stress response — no matter how well prepared."),
    ("🌱  Emotion regulation is a learnable skill.",
     "UC Berkeley's Greater Good Science Center confirms: children who learn to regulate emotional responses perform better academically, socially, and in long-term well-being."),
]
for i, (title, body) in enumerate(points):
    t = Inches(2.2) + i * Inches(1.5)
    add_textbox(s, Inches(0.65), t, Inches(11.5), Inches(0.42),
                title, font_size=14, bold=True, color=FOREST)
    add_textbox(s, Inches(0.65), t + Inches(0.42), Inches(11.5), Inches(0.72),
                body, font_size=12.5, color=MID)
bottom_line(s, "This is not a crisis intervention. It is an investment at exactly the right time.")
add_notes(s, """As engineers and scientists, you understand the concept of technical debt.

When you skip foundational work early in a project, you pay for it later — at much higher cost.

Emotional patterns work the same way. The anxious response, the shutdown, the self-doubt — if they are not addressed in childhood, they compound. They become the way this child responds to pressure at 16, at 25, at 40.

We are not here to fix a crisis. We are here to do the foundational work while it is easiest, most effective, and most lasting.""")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — The AAA Method
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
cream_bg(s)
header_bar(s, 9, 15, "What ThinkRoot Actually Does")
gold_headline(s, "Not therapy. Not tutoring. Not more content.\nA structured, repeatable method.")
add_textbox(s, Inches(0.65), Inches(2.2), Inches(11.5), Inches(0.4),
            "The AAA Method — three steps children learn to use in any difficult moment:",
            font_size=13, color=MID)
aaa_steps = [
    ("A  —  Awareness",
     "Notice what is happening in your mind and body before it escalates.\n\nCatch the thought, the feeling, the physical signal — early.\nThis is the step most children have never been taught."),
    ("A  —  Accept",
     "Acknowledge the thought or feeling without fighting it or pushing it away.\n\nSuppression makes patterns stronger. Acceptance is what allows transformation.\nThis step is what separates ThinkRoot from generic coping techniques."),
    ("A  —  Action",
     "Use a specific, practised technique to redirect and choose a steadier response.\n\nThe Five Power Technique is the concrete Action tool — under 3 minutes, no equipment, memorised and usable anywhere: before a test, after a mistake, at bedtime."),
]
step_w = Inches(3.9)
for i, (title, body) in enumerate(aaa_steps):
    l = Inches(0.45) + i * Inches(4.15)
    t = Inches(2.75)
    add_rect(s, l, t, step_w, Inches(3.3), FOREST)
    add_textbox(s, l + Inches(0.2), t + Inches(0.18), step_w - Inches(0.3), Inches(0.48),
                title, font_name="Playfair Display", font_size=17, bold=True, color=GOLD)
    add_textbox(s, l + Inches(0.2), t + Inches(0.72), step_w - Inches(0.3), Inches(2.4),
                body, font_size=11.5, color=WHITE)
add_rect(s, Inches(0.45), Inches(6.2), W - Inches(0.9), Inches(0.55), MIST)
add_textbox(s, Inches(0.65), Inches(6.27), W - Inches(1.2), Inches(0.45),
            "All three steps together create genuine inner agency — the ability to work with your own mind. We call this self-intelligence.",
            font_size=12.5, bold=True, color=FOREST, align=PP_ALIGN.CENTER)
add_notes(s, """So what does ThinkRoot actually do?

It is not therapy. We are not diagnosing or treating. It is an educational programme — practical, structured, and grounded in research.

The core method is three steps: Awareness, Accept, Action. Children learn to apply all three in sequence in any difficult moment.

The Accept step is the one most parents find surprising. Our instinct is to tell children "don't worry" or "stop thinking that." But that is suppression — and suppression makes patterns stronger over time. The Accept step teaches children to acknowledge what is happening without being ruled by it. That is what makes the transformation possible.

The Action step has a specific tool attached to it: the Five Power Technique. It is under three minutes, requires nothing, and can be used independently — before an exam, after a difficult conversation, at bedtime. Children memorise it so it is available the moment they need it.

Think of AAA as the operating system that runs underneath everything else your child does. Academics, sports, relationships, performance — all of these work better when the underlying system is stable.""")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — What Children Actually Do (5-Day Arc)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
cream_bg(s)
header_bar(s, 10, 15, "What Children Actually Do — The 5-Day Arc")
gold_headline(s, "Five days. Five building blocks. One complete toolkit.")
days = [
    ("Day 1", "The Brain Web",
     "How thoughts connect and spread. Children map their own thought patterns. Introduction to AAA and the Five Power Technique."),
    ("Day 2", "Brain Weather",
     "Emotional states mapped as weather — sun, storms, fog, rain. Three Specifics gratitude practice (UC Berkeley). Children learn to read their own inner forecast."),
    ("Day 3", "The Feeling Body",
     "Where feelings live in the body. Body scan and Heart Coherence Breathing (HeartMath Institute). Children learn to work with physical stress signals directly."),
    ("Day 4", "Upgrade the Thought",
     "Transforming storm thoughts into sunshine thoughts. Kindness Ripple activity. Self-compassion practice (Kristin Neff, UT Austin)."),
    ("Day 5", "I Am the Author",
     "Full integration. Children build their personal Resilience Toolkit. Five Power Technique led by children themselves. Completion certificates."),
]
dw = Inches(2.3)
dh = Inches(3.6)
dt = Inches(2.1)
for i, (day, title, desc) in enumerate(days):
    l = Inches(0.35) + i * Inches(2.55)
    add_rect(s, l, dt, dw, dh, FOREST)
    add_rect(s, l, dt, dw, Inches(0.52), GOLD)
    add_textbox(s, l + Inches(0.12), dt + Inches(0.06), dw - Inches(0.2), Inches(0.4),
                day, font_size=12, bold=True, color=FOREST, align=PP_ALIGN.CENTER)
    add_textbox(s, l + Inches(0.12), dt + Inches(0.6), dw - Inches(0.2), Inches(0.48),
                title, font_name="Playfair Display", font_size=13, bold=True, color=GOLD)
    add_textbox(s, l + Inches(0.12), dt + Inches(1.15), dw - Inches(0.2), dh - Inches(1.3),
                desc, font_size=10.5, color=WHITE)
add_rect(s, Inches(0.35), Inches(5.85), W - Inches(0.7), Inches(0.55), MIST)
add_textbox(s, Inches(0.55), Inches(5.92), W - Inches(1.0), Inches(0.45),
            "Every session: guided discussion · reflection · skill practice · home practice card  ·  Small groups · 1 hour · Grades 2–4 and 5–8 separated",
            font_size=11.5, color=FOREST, align=PP_ALIGN.CENTER)
add_notes(s, """Parents often ask me: what will my child actually do? Now I can show you.

[Walk through each day briefly — point to the card as you name it.]

Day 1 is about the brain and thoughts — children do a yarn activity to see how thoughts connect and spread. They learn the AAA method and the Five Power Technique on Day 1, because we want them practising it from the very first day.

Day 2 maps emotional states to weather. It sounds simple — and children love it — but the science behind it is UC Berkeley research on gratitude and positive states. Brain Weather gives children a language for what they are feeling without it feeling clinical.

Day 3 goes into the body. HeartMath Institute research shows that the heart and body respond to emotional states measurably. Children learn to use breath and awareness to shift a stress state physically, not just mentally.

Day 4 is where children actively transform their thought patterns — from storm thoughts to sunshine thoughts. And we introduce self-compassion, which research shows is the foundation for self-belief.

Day 5 is integration. Children build their own Resilience Toolkit — a personal one-page reference in their own words — and they lead each other through the Five Power Technique. They leave with something concrete they built themselves.

The 6-week format covers the same five topics with one week between each session, so children have time to practise in real life before the next concept arrives. It also adds a sixth session on building a personal long-term practice plan.""")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Results
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
cream_bg(s)
header_bar(s, 11, 15, "Real Results")
gold_headline(s, "From real children. In real situations.")
stories = [
    ('"She learned to pause before reacting. Her daily life became calmer. Her teachers noticed. Her relationships improved."',
     "— HM, Middle School"),
    ('"He used the technique during exams to relax, regulate, and stay composed. His performance improved."',
     "— RR, Middle School"),
    ('"She went from giving up when work felt difficult to breaking problems into steps and persisting."',
     "— Elementary school student"),
    ('"One student went from failing in math to improving his grades — and eventually helping classmates who were struggling in the same subject."',
     "— MT, Elementary School"),
]
sw = Inches(5.7)
sh = Inches(1.72)
for i, (quote, attr) in enumerate(stories):
    row = i // 2
    col = i % 2
    l = Inches(0.45) + col * Inches(6.95)
    t = Inches(2.1) + row * (sh + Inches(0.12))
    add_rect(s, l, t, sw, sh, MIST)
    add_rect(s, l, t, Inches(0.06), sh, GOLD)
    add_textbox(s, l + Inches(0.2), t + Inches(0.1), sw - Inches(0.3), sh - Inches(0.42),
                quote, font_size=11.5, color=FOREST)
    add_textbox(s, l + Inches(0.2), t + sh - Inches(0.38), sw - Inches(0.3), Inches(0.32),
                attr, font_size=11, bold=True, color=GOLD)
bottom_line(s, "These are not exceptional children. They are ordinary children who were given a tool they did not have before.")
add_notes(s, """I want to share a few real stories. These are children from this area — the same schools, the same communities your children are in.

[Read 2-3 of the stories naturally, don't rush.]

What strikes me about these outcomes is that none of them are about a dramatic transformation overnight. They are about a child who, in a difficult moment, now has something to reach for. A process. A skill.

That is all we are building. But that is everything.""")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Programme
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
cream_bg(s)
header_bar(s, 12, 15, "The Programme: Resilience Lab")
gold_headline(s, "Two formats. One goal. Right for Grades 2–8.")
# Left card
cw = Inches(5.7)
ch = Inches(3.6)
ct = Inches(2.1)
add_rect(s, Inches(0.45), ct, cw, ch, FOREST)
add_textbox(s, Inches(0.65), ct + Inches(0.18), cw - Inches(0.3), Inches(0.48),
            "5-Day Intensive", font_name="Playfair Display", font_size=20, bold=True, color=GOLD)
add_textbox_paras(s, Inches(0.65), ct + Inches(0.72), cw - Inches(0.3), ch - Inches(0.9),
    [("One focused week", False, 13, WHITE, PP_ALIGN.LEFT),
     ("1 hour per day, Monday–Friday", False, 13, WHITE, PP_ALIGN.LEFT),
     ("Small groups", False, 13, WHITE, PP_ALIGN.LEFT),
     ("", False, 12, WHITE, PP_ALIGN.LEFT),
     ("Best for: children who benefit from immersion and momentum", True, 12, GOLD, PP_ALIGN.LEFT)])
# Right card
add_rect(s, Inches(7.15), ct, cw, ch, FOREST)
add_textbox(s, Inches(7.35), ct + Inches(0.18), cw - Inches(0.3), Inches(0.48),
            "6-Week Series", font_name="Playfair Display", font_size=20, bold=True, color=GOLD)
add_textbox_paras(s, Inches(7.35), ct + Inches(0.72), cw - Inches(0.3), ch - Inches(0.9),
    [("One session per week for 6 weeks", False, 13, WHITE, PP_ALIGN.LEFT),
     ("1 hour per session", False, 13, WHITE, PP_ALIGN.LEFT),
     ("Small groups", False, 13, WHITE, PP_ALIGN.LEFT),
     ("", False, 12, WHITE, PP_ALIGN.LEFT),
     ("Best for: children who benefit from reflection time between sessions", True, 12, GOLD, PP_ALIGN.LEFT)])
# Both formats row
add_rect(s, Inches(0.45), ct + ch + Inches(0.12), W - Inches(0.9), Inches(0.75), MIST)
add_textbox(s, Inches(0.65), ct + ch + Inches(0.18), W - Inches(1.2), Inches(0.6),
            "Both formats: Grades 2–4 and Grades 5–8  ·  Mission San Jose, Fremont  ·  Online option on demand  ·  Summer 2026",
            font_size=12.5, color=FOREST, align=PP_ALIGN.CENTER)
add_notes(s, """We offer two formats to match different children and different schedules.

The content is the same. The difference is pace and how much time there is between sessions to practice.

Groups are small and formed by grade — so the conversations are appropriate and relevant to where each child actually is.

We are currently based in the Mission San Jose area of Fremont. If enough families in another part of the Bay Area are interested, we bring the programme to you.""")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — About Tina
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
cream_bg(s)
header_bar(s, 13, 15, "About Tina")
gold_headline(s, "I come from your world.")
# Left: photo placeholder
add_rect(s, Inches(0.45), Inches(2.1), Inches(3.8), Inches(4.5), MIST)
add_textbox(s, Inches(0.45), Inches(2.1), Inches(3.8), Inches(4.5),
            "[Photo of Tina]", font_size=14, color=MID, align=PP_ALIGN.CENTER)
# Right: credentials
creds = [
    "Tina Puri  |  Founder & Program Director",
    "",
    "✦  Certified Wellness Educator",
    "✦  Senior Peer Counselor — certified by the cities of Fremont, Union City, and Newark",
    "✦  Former software professional — 20+ years in the industry",
    "✦  Parent volunteer — Stanford University Challenge Success programme (Dr. Denise Pope)",
    "✦  Has worked with children, teens, adults, and seniors across schools and community settings",
    '✦  Co-founded Rewired to Inspire — four-year intergenerational community programme recognised\n    by Tri-City Voice and Congressman Michael M. Honda',
]
paras = []
for i, c in enumerate(creds):
    bold = i == 0
    sz = 15 if i == 0 else 12.5
    clr = FOREST if i == 0 else MID
    paras.append((c, bold, sz, clr, PP_ALIGN.LEFT))
add_textbox_paras(s, Inches(4.55), Inches(2.1), Inches(8.4), Inches(3.5), paras)
# Quote
add_rect(s, Inches(4.55), Inches(5.65), Inches(8.4), Inches(1.1), FOREST)
add_rect(s, Inches(4.55), Inches(5.65), Inches(0.06), Inches(1.1), GOLD)
add_textbox(s, Inches(4.75), Inches(5.75), Inches(8.1), Inches(0.95),
            '"I teach this because I have seen how deeply these skills can change a life. ThinkRoot is the result of years of personal practice, scientific study, and direct work with real students and families."',
            font_size=12, color=WHITE)
add_notes(s, """I want to take one minute to tell you why I built this.

I come from the same world many of you come from. I spent over two decades in software. I understand the environment your children are growing up in — the pressure, the competition, the expectation.

And I also know, from personal experience, that the skills I teach today — the ability to work with your own mind under pressure — are the ones that determine how everything else lands.

I built ThinkRoot because I looked around and saw that the children in this community had every advantage — except this one. And I knew what it could mean for them to have it early.""")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Closing
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, W, H, BARK)
add_rect(s, 0, 0, W, Inches(0.08), GOLD)
add_rect(s, 0, H - Inches(0.08), W, Inches(0.08), GOLD)
add_textbox(s, Inches(0.65), Inches(0.4), Inches(11.5), Inches(0.65),
            "You have invested in everything that goes in.",
            font_name="Playfair Display", font_size=24, bold=True, color=GOLD)
add_textbox(s, Inches(0.65), Inches(1.1), Inches(11.5), Inches(0.42),
            "What holds it all together?", font_name="Playfair Display", font_size=18,
            color=WHITE)
closing = [
    ("Your child has knowledge.", False, 15, MIST, PP_ALIGN.LEFT),
    ("But can they use it when the pressure is real?", True, 15, WHITE, PP_ALIGN.LEFT),
    ("", False, 8, WHITE, PP_ALIGN.LEFT),
    ("Your child has ability.", False, 15, MIST, PP_ALIGN.LEFT),
    ("But can they recover when they fall?", True, 15, WHITE, PP_ALIGN.LEFT),
    ("", False, 8, WHITE, PP_ALIGN.LEFT),
    ("Your child has opportunity.", False, 15, MIST, PP_ALIGN.LEFT),
    ("But can they believe in themselves when doubt rises?", True, 15, WHITE, PP_ALIGN.LEFT),
]
add_textbox_paras(s, Inches(0.65), Inches(1.7), Inches(11.5), Inches(3.8), closing)
add_rect(s, Inches(0.45), Inches(5.5), W - Inches(0.9), Inches(0.002), GOLD)
add_textbox(s, Inches(0.65), Inches(5.6), Inches(11.5), Inches(0.5),
            "The skills that answer yes to those questions — they are not born. They are built.",
            font_name="Playfair Display", font_size=16, bold=True, color=GOLD)
add_textbox(s, Inches(0.65), Inches(6.2), Inches(11.5), Inches(0.42),
            "And the window to build them well is now.",
            font_name="Playfair Display", font_size=15, color=WHITE)
add_notes(s, """I want to leave you with a thought.

We measure our investment in our children by what we give them — the tutoring, the programmes, the schools, the opportunities.

But what holds all of that together? What is the foundation on which everything else is built?

It is a child who knows how to handle difficulty. Who can recover. Who believes they are capable even when the evidence is not yet there.

That is not built by more content. It is built by practice. And it is built early.""")


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Next Step
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
cream_bg(s)
header_bar(s, 15, 15, "Next Step")
gold_headline(s, "No commitment. Just tell us you're interested.")
# CTA box
add_rect(s, Inches(0.45), Inches(2.1), W - Inches(0.9), Inches(0.62), FOREST)
add_textbox(s, Inches(0.65), Inches(2.2), W - Inches(1.2), Inches(0.45),
            "Join the Priority Interest List", font_name="Playfair Display",
            font_size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.65), Inches(2.85), W - Inches(1.2), Inches(0.55),
            "We are currently forming groups for Summer 2026.  Groups are small — formed by grade, location, and schedule.",
            font_size=13, color=MID, align=PP_ALIGN.CENTER)
steps_list = [
    ("1", "Visit the website and join the interest list"),
    ("2", "We will reach out when a group is forming near you"),
    ("3", "No payment. No obligation. Just priority access."),
]
for i, (num, txt) in enumerate(steps_list):
    t = Inches(3.55) + i * Inches(0.72)
    add_rect(s, Inches(2.0), t, Inches(0.5), Inches(0.5), GOLD)
    add_textbox(s, Inches(2.0), t, Inches(0.5), Inches(0.5),
                num, font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(2.65), t + Inches(0.06), Inches(8.5), Inches(0.42),
                txt, font_size=14, color=FOREST)
add_rect(s, Inches(0.45), Inches(5.85), W - Inches(0.9), Inches(0.001), GOLD)
contact_paras = [
    ("🌐  vjra.net", False, 13, FOREST, PP_ALIGN.LEFT),
    ("📧  thinkyoungleaders@gmail.com", False, 13, FOREST, PP_ALIGN.LEFT),
]
add_textbox_paras(s, Inches(2.0), Inches(5.95), Inches(9), Inches(0.75), contact_paras)
add_textbox(s, Inches(0.65), Inches(6.75), W - Inches(1.2), Inches(0.45),
            '"The earlier you express interest, the better we can match your child to the right group."',
            font_size=12.5, color=MID, align=PP_ALIGN.CENTER)
add_notes(s, """The next step is simple and risk-free.

All I am asking today is that if this resonated — if you recognise your child in anything we talked about — you join the interest list. That is it. No payment. No commitment. Just your name and a few basics so we can match your child to the right group when it forms.

Groups are small by design. That is what makes the work effective. And because of that, they fill based on interest and timing.

I am happy to answer any questions now. And please feel free to reach out directly after today.

Thank you for your time and for caring enough to be here.""")


# ── Save ─────────────────────────────────────────────────────────────────────
out_path = "/Users/tina/Documents/code/thinkroot-site/launch-assets/presentation-2026-04-27.pptx"
prs.save(out_path)
print(f"Saved → {out_path}")
